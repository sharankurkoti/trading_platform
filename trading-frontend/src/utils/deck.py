from pptx import Presentation

# Create presentation
prs = Presentation()

slides_content = [
    {
        "layout": 0,
        "title": "DIT: Digitisation of International Trading",
        "subtitle": "The Integrated Future of Global Trade Finance\n\nPresenter: [Your Name], [Your Title]\nContact: [Your Email] | [Your Website]"
    },
    {
        "layout": 1,
        "title": "The $9 Trillion Bottleneck: Why Global Trade is Broken",
        "content": (
            "✋ Manual & Inefficient: Paper-heavy processes (e.g., Letters of Credit) create delays of weeks and costly errors.\n\n"
            "👁️ Lack of Transparency: Limited visibility leads to disputes, fraud, and a lack of trust among parties.\n\n"
            "💸 High Costs & Risks: Significant currency volatility and financing barriers stifle growth, especially for SMEs."
        )
    },
    {
        "layout": 1,
        "title": "Our Vision: One Platform, One Workflow, Zero Paper",
        "content": (
            "DIT is a unified, blockchain-enabled ecosystem that digitizes and connects the entire trade finance lifecycle—from "
            "agreement to settlement—creating a new standard for global trade."
        )
    },
    {
        "layout": 1,
        "title": "Our Integrated Service Suite",
        "content": (
            "1. Digital LC Management: End-to-end digital creation, automated compliance checks, and a secure document hub.\n\n"
            "2. Trade Financing Marketplace: A dynamic platform connecting sellers with funders, using digital LCs as verifiable collateral.\n\n"
            "3. Intelligent Trade Exchange: A neutral ground for negotiation, smart contract generation, and integrated currency risk management."
        )
    },
    {
        "layout": 1,
        "title": "The Seamless Digital Trade Journey",
        "content": (
            "Agree: Buyer and Seller agree on terms on the Trade Exchange.\n\n"
            "Secure: A Digital LC is issued and advised on the LC Management platform in minutes.\n\n"
            "Finance: Seller requests immediate pre-shipment funding from the Financing Marketplace.\n\n"
            "Verify: Goods are shipped; digital documents are uploaded and automatically verified by AI.\n\n"
            "Settle: Payment is released automatically; currencies are exchanged at pre-agreed rates."
        )
    },
    {
        "layout": 1,
        "title": "Service 1: Digital LC Management - Automating Trust & Compliance",
        "content": (
            "Digital LC creation with smart, compliant templates.\n\n"
            "AI-powered document discrepancy detection.\n\n"
            "Immutable audit trail via blockchain.\n\n"
            "Secure, role-based messaging hub for all parties.\n\n"
            "Tech Stack: Python, React, PostgreSQL, AWS S3, Blockchain."
        )
    },
    {
        "layout": 1,
        "title": "Service 2: Trade Financing Marketplace - Unlocking Liquidity",
        "content": (
            "Connects sellers with a network of banks and institutional funders.\n\n"
            "Uses the digital LC as secure, verifiable collateral, de-risking lenders.\n\n"
            "Features a risk dashboard with AI-driven credit assessment.\n\n"
            "Automated repayment upon transaction completion.\n\n"
            "Tech Stack: Node.js, React, PostgreSQL, Blockchain."
        )
    },
    {
        "layout": 1,
        "title": "Service 3: The Trade Exchange - Where Global Trade Negotiates",
        "content": (
            "Facilitates trade matching and negotiation.\n\n"
            "Generates smart, digital sales contracts.\n\n"
            "Integrated real-time forex API with hedging tools to mitigate currency risk.\n\n"
            "Direct gateway to initiate LCs and request financing.\n\n"
            "Tech Stack: Python/Ruby, Vue.js, PostgreSQL, Forex APIs."
        )
    },
    {
        "layout": 1,
        "title": "Built for Security, Scale, & Interoperability",
        "content": (
            "Architecture: Decentralized, microservices-based for resilience.\n\n"
            "APIs & Webhooks: Ensure seamless, real-time data flow between services.\n\n"
            "Security: OAuth 2.0, Role-Based Access Control (RBAC), and Blockchain immutability.\n\n"
            "Stack: Python, Java, Node.js, React, Vue, PostgreSQL, AWS."
        )
    },
    {
        "layout": 1,
        "title": "Why We Win: The End-to-End Advantage",
        "content": (
            "Full-Stack Solution: Only DIT integrates the entire trade process. Competitors are point solutions.\n\n"
            "Frictionless Experience: Dramatically reduces time and cost from weeks to days.\n\n"
            "Network Effects: More users → more data → better risk AI → more liquidity → more users."
        )
    },
    {
        "layout": 1,
        "title": "Capturing a Share of a Trillion-Dollar Market",
        "content": (
            "Global trade finance market is worth $9+ trillion annually.\n\n"
            "Target Customers: SMEs, Corporates, Banks, Financiers.\n\n"
            "Revenue Streams: Transaction fees, financing facilitation fees, FX spread."
        )
    },
    {
        "layout": 1,
        "title": "Our Path to Market Leadership",
        "content": (
            "Phase 1 (Launch - Now): LC Management MVP with pilot bank partners. (Focus on Proven Need)\n\n"
            "Phase 2 (Scale - Next 12 months): Integrate & launch Trade Financing Marketplace. (Unlock New Revenue)\n\n"
            "Phase 3 (Expand - 18+ months): Roll out full Trade Exchange with advanced FX capabilities. (Become the Ecosystem)"
        )
    },
    {
        "layout": 1,
        "title": "The Right Team to Execute",
        "content": (
            "[Your Name], CEO: [e.g., 10+ years in trade finance and FinTech leadership].\n\n"
            "[Team Member Name], CTO: [e.g., Expert in building scalable SaaS platforms and blockchain].\n\n"
            "[Advisor Name], Advisor: [e.g., Former MD of [Major Bank], deep industry connections]."
        )
    },
    {
        "layout": 1,
        "title": "We Are Seeking [$X] in Seed Funding",
        "content": (
            "40% - Product Development: Scaling the engineering and AI teams.\n\n"
            "30% - Commercialization: Sales, marketing, and establishing pilot programs.\n\n"
            "20% - Operations: Key hires, legal, and administrative costs.\n\n"
            "10% - Contingency."
        )
    },
    {
        "layout": 0,
        "title": "Thank You",
        "subtitle": (
            "Let's Build the Future of Trade Together\n\nDIT - Digitisation of International Trading\n\n"
            "[Your Name]\n[Your Email]\n[Your Website/LinkedIn]"
        )
    }
]

# Add slides
for slide_data in slides_content:
    slide_layout = prs.slide_layouts[slide_data["layout"]]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = slide_data["title"]
    
    if slide_data["layout"] == 0:
        subtitle = slide.placeholders[1]
        subtitle.text = slide_data["subtitle"]
    else:
        content = slide.placeholders[1]
        content.text = slide_data["content"]

# Save presentation
prs.save("DIT_Digitisation_of_International_Trading_PitchDeck.pptx")
print("Presentation saved as DIT_Digitisation_of_International_Trading_PitchDeck.pptx")